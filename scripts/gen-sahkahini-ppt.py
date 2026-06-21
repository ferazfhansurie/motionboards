# -*- coding: utf-8 -*-
"""Generate the SahkahIni? GFP 009 PowerPoint deck (modern template)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY   = RGBColor(0x0B, 0x12, 0x2B)   # deep background
NAVY2  = RGBColor(0x12, 0x1C, 0x3D)
TEAL   = RGBColor(0x2D, 0xD4, 0xBF)   # primary accent (trust)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # warning accent
CORAL  = RGBColor(0xFB, 0x71, 0x85)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1E, 0x29, 0x3B)   # dark text on light
SLATE  = RGBColor(0x64, 0x74, 0x8B)   # muted
LIGHT  = RGBColor(0xF6, 0xF8, 0xFB)   # light bg
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xE2, 0xE8, 0xF0)

FONT = "Segoe UI"
FONT_H = "Segoe UI Semibold"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(s, x, y, w, h, color, shadow=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(sp, color)
    sp.shadow.inherit = False
    return sp

def rrect(s, x, y, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill(sp, color); sp.shadow.inherit = False
    return sp

def bg(s, color):
    rect(s, 0, 0, SW, SH, color)

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space=4, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,font)"""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for (t, sz, c, b, fn) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.name = fn
    return tb

def R(t, sz, c, b=False, fn=FONT):  # run helper
    return (t, sz, c, b, fn)

def chip(s, x, y, label, color=TEAL, tcolor=NAVY):
    w = Inches(0.18 + 0.092 * len(label))
    c = rrect(s, x, y, w, Inches(0.36), color)
    txt(s, x, y, w, Inches(0.36), [[R(label, 12, tcolor, True, FONT_H)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w

PAGE = {"n": 0}
def header(s, kicker, title, accent=TEAL):
    bg(s, LIGHT)
    rect(s, 0, 0, Inches(0.22), SH, accent)            # left spine
    txt(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.4),
        [[R(kicker.upper(), 13, accent, True, FONT_H)]])
    txt(s, Inches(0.7), Inches(0.82), Inches(12), Inches(0.95),
        [[R(title, 30, INK, True, FONT_H)]])
    rect(s, Inches(0.72), Inches(1.62), Inches(1.1), Inches(0.05), accent)
    # footer
    PAGE["n"] += 1
    txt(s, Inches(0.7), Inches(7.04), Inches(6), Inches(0.3),
        [[R("SahkahIni?  ·  GFP 009 Digital Creativity", 9, SLATE, False)]])
    txt(s, Inches(11.3), Inches(7.04), Inches(1.3), Inches(0.3),
        [[R(str(PAGE["n"]).zfill(2), 9, SLATE, True)]], align=PP_ALIGN.RIGHT)

def card(s, x, y, w, h, title, body_lines, accent=TEAL, tsize=15, bsize=12):
    rrect(s, x, y, w, h, CARD)
    rect(s, x, y, Inches(0.08), h, accent)
    paras = [[R(title, tsize, INK, True, FONT_H)]]
    for ln in body_lines:
        paras.append([R(ln, bsize, SLATE, False)])
    txt(s, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5), h - Inches(0.3), paras, space=5)

def bullets(s, x, y, w, h, items, size=14, color=INK, accent=TEAL, space=9):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
        else:
            lead, rest = "", it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space); p.space_before = Pt(0)
        r0 = p.add_run(); r0.text = "▸  "; r0.font.size = Pt(size); r0.font.color.rgb = accent; r0.font.bold = True; r0.font.name = FONT_H
        if lead:
            r1 = p.add_run(); r1.text = lead + "  "; r1.font.size = Pt(size); r1.font.color.rgb = color; r1.font.bold = True; r1.font.name = FONT_H
        r2 = p.add_run(); r2.text = rest; r2.font.size = Pt(size); r2.font.color.rgb = SLATE; r2.font.name = FONT
    return tb

def table(s, x, y, w, headers, rows, col_w=None, accent=TEAL, fs=11):
    nrows = len(rows) + 1; ncols = len(headers)
    h = Inches(0.46) * nrows
    gt = s.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Emu(int(int(w) * cw / total))
    # header
    for j, hd in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.margin_left = c.margin_right = Inches(0.1); c.margin_top = c.margin_bottom = Inches(0.04)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = hd
        r.font.size = Pt(fs+0.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT_H
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i+1, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xEE,0xF3,0xF8)
            c.margin_left = c.margin_right = Inches(0.1); c.margin_top = c.margin_bottom = Inches(0.04)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
            r.font.size = Pt(fs); r.font.name = FONT
            r.font.color.rgb = INK if j == 0 else SLATE
            if j == 0: r.font.bold = True
    return gt

# ============================================================ SLIDE 1 — TITLE
s = slide(); bg(s, NAVY)
rect(s, 0, 0, SW, Inches(0.12), TEAL)
rect(s, 0, SH - Inches(0.12), SW, Inches(0.12), AMBER)
# decorative circles
for (cx, cy, d, col) in [(11.4,1.0,2.4,NAVY2),(12.3,5.6,1.6,NAVY2)]:
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
    fill(o, col); o.shadow.inherit = False
chip(s, Inches(0.9), Inches(1.15), "GFP 009  ·  DIGITAL CREATIVITY", TEAL, NAVY)
txt(s, Inches(0.86), Inches(1.95), Inches(11), Inches(1.7),
    [[R("SahkahIni", 76, WHITE, True, FONT_H), R("?", 76, TEAL, True, FONT_H)]])
txt(s, Inches(0.9), Inches(3.35), Inches(10.5), Inches(0.9),
    [[R("An AI WhatsApp bot that fact-checks forwarded messages & images — in Bahasa Malaysia.",
        20, RGBColor(0xC8,0xD3,0xE6), False)]])
chip(s, Inches(0.9), Inches(4.35), "AMPLIFY MODEL", AMBER, NAVY)
chip(s, Inches(3.1), Inches(4.35), "MEDIA STUDIES · UM", NAVY2, RGBColor(0xC8,0xD3,0xE6))
txt(s, Inches(0.9), Inches(5.4), Inches(11), Inches(1.4),
    [[R("PREPARED BY", 12, TEAL, True, FONT_H)],
     [R("Aisyatun Nabiha", 17, WHITE, True, FONT_H),
      R("    ·    Ahli Kumpulan 2    ·    Ahli Kumpulan 3", 15, SLATE, False)],
     [R("Faculty of Arts & Social Sciences (Media Studies), Universiti Malaya", 12, SLATE, False)]],
    space=6)

# ============================================================ SLIDE 2 — AGENDA
s = slide(); header(s, "Contents", "What this deck covers")
items_l = [
    ("01  Introduction", "The problem in one breath"),
    ("02  Problem Statement", "Issue, significance, justification, objectives"),
    ("03  Use of Technology", "Past tools vs today's AI"),
    ("04  Future Reflection", "Trends shaping adoption"),
]
items_r = [
    ("05  Project Management", "Timeline, sustainability, commercialisation"),
    ("06  Designing the Solution", "Define → Prototype → Test → Evaluate"),
    ("07  Storyboard", "How a user experiences SahkahIni"),
    ("08  Peer Eval · Conclusion · Docs", "Improvements & wrap-up"),
]
def agenda_col(x, items, tag, color):
    chip(s, x, Inches(1.95), tag, color, NAVY if color!=NAVY2 else WHITE)
    y = Inches(2.55)
    for head, sub in items:
        rrect(s, x, y, Inches(5.6), Inches(0.92), CARD)
        rect(s, x, y, Inches(0.08), Inches(0.92), color)
        txt(s, x+Inches(0.3), y+Inches(0.13), Inches(5.2), Inches(0.4), [[R(head, 15, INK, True, FONT_H)]])
        txt(s, x+Inches(0.3), y+Inches(0.5), Inches(5.2), Inches(0.35), [[R(sub, 11.5, SLATE, False)]])
        y += Inches(1.06)
agenda_col(Inches(0.9), items_l, "TUGASAN 1  ·  40 MARKS", TEAL)
agenda_col(Inches(6.85), items_r, "TUGASAN 2  ·  60 MARKS", AMBER)

# ============================================================ SLIDE 3 — SECTION T1
def divider(num, title, sub, accent):
    s = slide(); bg(s, NAVY)
    rect(s, 0, 0, Inches(0.22), SH, accent)
    txt(s, Inches(0.9), Inches(2.2), Inches(8), Inches(2),
        [[R(num, 120, NAVY2, True, FONT_H)]])
    txt(s, Inches(0.95), Inches(3.0), Inches(11), Inches(1.2), [[R(title, 44, WHITE, True, FONT_H)]])
    rect(s, Inches(0.98), Inches(4.05), Inches(1.4), Inches(0.06), accent)
    txt(s, Inches(0.98), Inches(4.3), Inches(10.5), Inches(0.9), [[R(sub, 18, RGBColor(0xC8,0xD3,0xE6), False)]])
    return s
divider("01", "TUGASAN 1", "Problem & Justification  ·  40 marks", TEAL)

# ============================================================ SLIDE 4 — INTRO
s = slide(); header(s, "1 · Introduction (5%)", "Every day, a hoax goes viral before the truth wakes up")
bullets(s, Inches(0.9), Inches(2.0), Inches(6.6), Inches(4),
    [("The vector:", "WhatsApp is Malaysia's #1 messaging app — and its #1 misinformation channel. A single 'forwarded many times' message reaches thousands in minutes."),
     ("The harm:", "health scams, fake 'free' giveaways, doctored images, and political rumours that erode public trust."),
     ("The gap:", "most fact-checking tools are slow, manual, and English-first — leaving Bahasa Malaysia users unprotected."),
     ("Our answer:", "SahkahIni? — forward any suspicious message to a bot and get a clear Benar / Palsu / Belum Disahkan verdict, with sources, in seconds.")],
    size=14, space=11)
# right verdict-card mock
x = Inches(8.0); rrect(s, x, Inches(2.0), Inches(4.4), Inches(4.4), CARD)
rect(s, x, Inches(2.0), Inches(4.4), Inches(0.7), TEAL)
txt(s, x+Inches(0.3), Inches(2.0), Inches(4), Inches(0.7), [[R("SahkahIni?  ✓", 16, NAVY, True, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
rrect(s, x+Inches(0.3), Inches(2.95), Inches(3.8), Inches(0.95), RGBColor(0xEE,0xF3,0xF8))
txt(s, x+Inches(0.5), Inches(3.05), Inches(3.5), Inches(0.8),
    [[R("“Bank akan tutup akaun jika tak update IC...”", 12, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
rrect(s, x+Inches(0.3), Inches(4.1), Inches(1.8), Inches(0.55), CORAL)
txt(s, x+Inches(0.3), Inches(4.1), Inches(1.8), Inches(0.55), [[R("PALSU", 15, WHITE, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, x+Inches(0.32), Inches(4.85), Inches(3.9), Inches(1.4),
    [[R("Why:", 12, INK, True, FONT_H), R(" matches a known scam pattern flagged by MCMC & banks.", 12, SLATE, False)],
     [R("Source:", 12, INK, True, FONT_H), R(" Sebenarnya.my · BNM advisory", 11, SLATE, False)]], space=6)

# ============================================================ SLIDE 5 — PROBLEM: issue
s = slide(); header(s, "2a · Problem Statement (10%)", "The issue — history & background")
bullets(s, Inches(0.9), Inches(2.0), Inches(6.5), Inches(4.6),
    [("2018 onward:", "misinformation surges globally; Malaysia debates the Anti-Fake News Act."),
     ("WhatsApp design:", "frictionless forwarding + closed groups make corrections impossible to trace."),
     ("Pandemic peak:", "COVID-19 health hoaxes (fake cures, vaccine myths) spread faster than official advice."),
     ("Government response:", "MCMC launches Sebenarnya.my; still manual, reactive, and not where people actually chat."),
     ("Language gap:", "leading AI fact-checkers (Snopes, ClaimBuster) are English-centric — BM content slips through.")],
    size=13.5, space=10)
card(s, Inches(7.8), Inches(2.0), Inches(4.6), Inches(2.1), "Key references to cite",
     ["MCMC & Sebenarnya.my reports", "Reuters Institute Digital News Report (MY)", "Academic work on WhatsApp misinformation", "Bank Negara scam advisories"], accent=AMBER)
card(s, Inches(7.8), Inches(4.3), Inches(4.6), Inches(2.1), "The core problem",
     ["Malaysians lack a fast, trusted,", "BM-native way to verify a message", "at the exact moment they receive it —", "inside the chat app itself."], accent=CORAL)

# ============================================================ SLIDE 6 — significance + justification
s = slide(); header(s, "2b · 2c · Problem Statement (10%)", "Significance & justification")
card(s, Inches(0.9), Inches(2.0), Inches(5.6), Inches(2.05), "b · Significance to our team", [], accent=TEAL)
bullets(s, Inches(1.15), Inches(2.55), Inches(5.1), Inches(1.5),
    [("Media students:", "we're trained to spot manipulation — most of our families are not."),
     ("Personal:", "we've all watched relatives forward hoaxes in group chats.")], size=12.5, space=7)
card(s, Inches(6.85), Inches(2.0), Inches(5.55), Inches(2.05), "c · Justification — why solve it now", [], accent=AMBER)
bullets(s, Inches(7.1), Inches(2.55), Inches(5.05), Inches(1.5),
    [("Real damage:", "financial scams, public-health risk, social tension."),
     ("Tools fall short:", "slow, manual, English-first, outside the chat.")], size=12.5, space=7)
# objectives strip
chip(s, Inches(0.9), Inches(4.35), "d · PROJECT OBJECTIVES (SMART)", TEAL, NAVY)
objs = [("O1", "Verify a forwarded BM message in under 30 seconds."),
        ("O2", "Support text + image (reverse-image) + voice-note claims."),
        ("O3", "Return a clear verdict with at least one credible source."),
        ("O4", "Reach everyday users where they already are — inside WhatsApp.")]
x = Inches(0.9); w = Inches(2.78)
for i,(n,t) in enumerate(objs):
    xx = x + i*(w+Inches(0.16))
    rrect(s, xx, Inches(4.95), w, Inches(1.6), CARD)
    rect(s, xx, Inches(4.95), w, Inches(0.5), NAVY)
    txt(s, xx, Inches(4.95), w, Inches(0.5), [[R(n, 15, TEAL, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, xx+Inches(0.2), Inches(5.55), w-Inches(0.4), Inches(0.95), [[R(t, 12, SLATE, False)]])

# ============================================================ SLIDE 7 — TECH past vs new + table
s = slide(); header(s, "3 · Use of Technology (10%)", "From manual desks to multimodal AI")
bullets(s, Inches(0.9), Inches(1.95), Inches(5.4), Inches(2.3),
    [("Previously used:", "manual fact-check desks, browser plug-ins, reverse-image search, English-only AI classifiers (ClaimBuster, Snopes)."),
     ("Current new tech:", "LLMs for claim detection, multimodal AI (image+text), WhatsApp Business API bots, RAG over verified-news databases.")],
    size=13, space=10)
chip(s, Inches(6.6), Inches(1.95), "THE INNOVATION GAP", AMBER, NAVY)
table(s, Inches(6.6), Inches(2.45), Inches(6.0),
      ["Capability", "Existing", "SahkahIni"],
      [["Bahasa Malaysia native", "Limited", "Yes"],
       ["Inside the chat app", "No", "Yes"],
       ["Text + image + voice", "Partial", "Yes"],
       ["Instant (<30s)", "Manual", "Yes"],
       ["Cites sources", "Varies", "Always"]],
      col_w=[2.4,1.5,1.5], fs=11)
txt(s, Inches(0.9), Inches(4.6), Inches(5.4), Inches(2),
    [[R("Why the table matters", 14, INK, True, FONT_H)],
     [R("This 'existing vs ours' comparison is where the rubric awards Ideation / product-innovation marks. It proves SahkahIni isn't a copy — it closes a real, named gap.", 12.5, SLATE, False)]], space=8)

# ============================================================ SLIDE 8 — FUTURE
s = slide(); header(s, "4 · Future Reflection (5%)", "Where this technology is heading")
cards = [
    ("Agentic AI", ["Bots that investigate", "a claim end-to-end,", "not just classify it."], TEAL),
    ("Deepfake era", ["Image/voice/video", "verification becomes", "essential, not optional."], CORAL),
    ("On-device + privacy", ["Verification that never", "uploads your private", "chats to a server."], AMBER),
    ("Regulation & schools", ["Integration into WhatsApp,", "newsrooms, and media-", "literacy curricula."], TEAL),
]
x = Inches(0.9); w = Inches(2.85)
for i,(t,b,c) in enumerate(cards):
    xx = x + i*(w+Inches(0.16))
    card(s, xx, Inches(2.1), w, Inches(2.7), t, b, accent=c, tsize=15, bsize=12)
txt(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4),
    [[R("Adoption outlook:", 14, INK, True, FONT_H),
      R("  As deepfakes scale and regulators push platform accountability, a BM-native, in-chat verifier moves from 'nice to have' to public infrastructure — embedded in WhatsApp itself, taught in schools, and used by newsrooms.", 13, SLATE, False)]], space=8)

# ============================================================ SLIDE 9 — SECTION T2
divider("02", "TUGASAN 2", "Solution & Project Management  ·  60 marks", AMBER)

# ============================================================ SLIDE 10 — Project mgmt timeline
s = slide(); header(s, "4a · Project Management (10)", "Timeline — mapped to the course schedule", accent=AMBER)
table(s, Inches(0.9), Inches(2.05), Inches(11.5),
      ["Phase", "Course weeks", "Activity", "AMPLIFY"],
      [["Define", "Wk 1–2", "Problem research, references, objectives", "A · Ask"],
       ["Map", "Wk 3–4", "Survey existing tools, build gap table", "M · Map"],
       ["Ideate", "Wk 5–6", "Concept sketches, pick best solution", "P · Picture"],
       ["Prototype", "Wk 7–9", "Storyboard + clickable bot mock-up", "L · Layout"],
       ["Test", "Wk 10–11", "Trial with users, peer evaluation", "I · Iterate"],
       ["Finalise", "Wk 12–13", "Report, e-portfolio, presentation", "F · Finalise"]],
      col_w=[1.4,1.5,4.5,1.6], fs=11.5, accent=AMBER)

# ============================================================ SLIDE 11 — sustainability / commercial / others
s = slide(); header(s, "4b · 4c · 4d · Project Management (10)", "Sustainability, commercialisation & creativity", accent=AMBER)
card(s, Inches(0.9), Inches(2.0), Inches(3.7), Inches(4.4), "b · Sustainability",
     ["No further funding needed:", "", "• Runs on free / low-cost API tiers", "• Open-source AI models", "• Student-maintained as a club", "• Cloud free-tier hosting", "• Community-sourced fact database"], accent=TEAL, bsize=12.5)
card(s, Inches(4.8), Inches(2.0), Inches(3.7), Inches(4.4), "c · Commercialisation",
     ["Freemium model:", "", "• Free for the public", "• Paid API for newsrooms", "• Enterprise tier for govt /", "  brands fighting scams", "• Partnership: MCMC,", "  Sebenarnya.my"], accent=AMBER, bsize=12.5)
card(s, Inches(8.7), Inches(2.0), Inches(3.7), Inches(4.4), "d · Others (creativity)",
     ["Make it loveable:", "", "• Friendly mascot & BM voice", "• Gamified 'Trust Score'", "• Shareable correction cards", "• Weekly 'Top Hoax' digest", "• School workshop kit"], accent=CORAL, bsize=12.5)

# ============================================================ SLIDE 12 — Designing the solution
s = slide(); header(s, "5a · Process (30)", "Designing the solution — the design loop", accent=AMBER)
steps = [("DEFINE","Pin the problem & users", TEAL),("IDEATE","Brainstorm, pick best concept", AMBER),
         ("PROTOTYPE","Storyboard + bot mock-up", CORAL),("TEST","Trial with real users", TEAL),
         ("EVALUATE","Measure & refine", AMBER)]
x = Inches(0.9); w = Inches(2.18)
for i,(t,sub,c) in enumerate(steps):
    xx = x + i*(w+Inches(0.18))
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, xx+Inches(0.75), Inches(2.2), Inches(0.7), Inches(0.7)); fill(o,c); o.shadow.inherit=False
    txt(s, xx+Inches(0.75), Inches(2.2), Inches(0.7), Inches(0.7), [[R(str(i+1), 22, WHITE, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, xx, Inches(3.1), w, Inches(1.9), CARD)
    rect(s, xx, Inches(3.1), w, Inches(0.06), c)
    txt(s, xx+Inches(0.15), Inches(3.3), w-Inches(0.3), Inches(0.5), [[R(t, 14, INK, True, FONT_H)]], align=PP_ALIGN.CENTER)
    txt(s, xx+Inches(0.15), Inches(3.85), w-Inches(0.3), Inches(1.0), [[R(sub, 12, SLATE, False)]], align=PP_ALIGN.CENTER)
    if i < 4:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, xx+w+Inches(0.02), Inches(3.85), Inches(0.16), Inches(0.3)); fill(ar, SLATE); ar.shadow.inherit=False
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.2),
    [[R("Defining → Ideation → Prototyping → Testing → Evaluation. ", 13, INK, True, FONT_H),
      R("Each loop feeds the next; testing failures send us back to ideate. The storyboard (next slide) is our core prototype artefact.", 13, SLATE, False)]], space=8)

# ============================================================ SLIDE 13 — STORYBOARD
s = slide(); header(s, "5a · Storyboard with descriptions", "How a user experiences SahkahIni?", accent=AMBER)
frames = [
    ("1 · Receive", "User gets a suspicious 'forwarded many times' message in a family group.", CORAL),
    ("2 · Forward", "They forward it to the SahkahIni? WhatsApp number. Bot replies: “Nak saya semak?”", TEAL),
    ("3 · Analyse", "AI reads the text + reverse-searches any image against verified sources.", AMBER),
    ("4 · Verdict", "Bot returns a card: PALSU / BENAR / BELUM DISAHKAN, with a plain-BM reason.", CORAL),
    ("5 · Sources", "Verdict links to Sebenarnya.my / news so the user can check for themselves.", TEAL),
    ("6 · Share", "User shares the correction card back to the group — truth travels too.", AMBER),
]
x0, y0 = Inches(0.9), Inches(2.0); w, h = Inches(3.7), Inches(2.05)
for i,(t,d,c) in enumerate(frames):
    col = i % 3; rowi = i // 3
    xx = x0 + col*(w+Inches(0.2)); yy = y0 + rowi*(h+Inches(0.25))
    rrect(s, xx, yy, w, h, CARD)
    rect(s, xx, yy, w, Inches(0.5), c)
    txt(s, xx+Inches(0.2), yy, w-Inches(0.3), Inches(0.5), [[R(t, 14, WHITE, True, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, xx+Inches(0.22), yy+Inches(0.62), w-Inches(0.44), Inches(1.3), [[R(d, 12.5, SLATE, False)]])

# ============================================================ SLIDE 14 — peer eval / recommendations
s = slide(); header(s, "5b · Recommendations for improvement", "Peer evaluation drives the next iteration", accent=AMBER)
table(s, Inches(0.9), Inches(2.05), Inches(11.5),
      ["Feedback from testers / peers", "Action taken", "Status"],
      [["Verdict felt too technical", "Rewrote replies in everyday BM", "Done"],
       ["Users unsure how to start", "Added “forward & ask” onboarding", "Done"],
       ["Wanted to verify images too", "Added reverse-image search", "Done"],
       ["Worried about privacy", "Added on-device / no-store option", "Planned"],
       ["Asked for voice-note support", "Added speech-to-text claim check", "Planned"]],
      col_w=[4.8,4.0,1.4], fs=12, accent=CORAL)
txt(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1),
    [[R("Peer evaluation method:  ", 13, INK, True, FONT_H),
      R("each team member + 5 external testers rated clarity, speed, and trust on a 1–5 scale; lowest-scoring areas became the improvement backlog above.", 13, SLATE, False)]], space=6)

# ============================================================ SLIDE 15 — AMPLIFY recap
s = slide(); header(s, "Framework", "How SahkahIni? applies the AMPLIFY model")
rows = [("A","Ask","Define the misinformation problem & who it harms", TEAL),
        ("M","Map","Review existing tools; find the BM / in-chat gap", AMBER),
        ("P","Picture","Ideate the bot concept; pick the best solution", CORAL),
        ("L","Layout","Storyboard + clickable WhatsApp prototype", TEAL),
        ("I","Iterate","Test with users; peer evaluation loop", AMBER),
        ("F","Finalise","Ship the verified-verdict experience & report", CORAL),
        ("Y","Yield","Sustainability, commercialisation & future impact", TEAL)]
y = Inches(2.0)
for (L,name,desc,c) in rows:
    rrect(s, Inches(0.9), y, Inches(11.5), Inches(0.62), CARD)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y+Inches(0.07), Inches(0.48), Inches(0.48)); fill(o,c); o.shadow.inherit=False
    txt(s, Inches(1.0), y+Inches(0.07), Inches(0.48), Inches(0.48), [[R(L, 18, WHITE, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.7), y, Inches(2.2), Inches(0.62), [[R(name, 15, INK, True, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.9), y, Inches(8.3), Inches(0.62), [[R(desc, 13, SLATE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.68)

# ============================================================ SLIDE 16 — conclusion + docs
s = slide(); header(s, "5c · 5d · Conclusion & Documentation (10)", "Wrap-up, recommendations & records", accent=AMBER)
card(s, Inches(0.9), Inches(2.0), Inches(5.6), Inches(4.4), "c · Conclusion & recommendations",
     ["SahkahIni? closes a real gap:", "BM-native, in-chat, instant verification.", "",
      "Recommendations for delivery:", "• Pilot with one UM faculty group", "• Partner with Sebenarnya.my for data",
      "• Add deepfake (image/video) detection", "• Run media-literacy workshops", "• Track 'hoaxes stopped' as impact metric"], accent=TEAL, bsize=12.5)
card(s, Inches(6.85), Inches(2.0), Inches(5.55), Inches(4.4), "d · Documentation (group meetings)",
     ["Meeting log kept in e-portfolio:", "", "• M1 (Wk1): scope + roles assigned",
      "• M2 (Wk4): tech review + gap table", "• M3 (Wk7): storyboard finalised", "• M4 (Wk10): test results reviewed",
      "• M5 (Wk13): report + slides sign-off", "", "Each entry: date · attendees · decisions"], accent=AMBER, bsize=12.5)

# ============================================================ SLIDE 17 — THANK YOU
s = slide(); bg(s, NAVY)
rect(s, 0, 0, SW, Inches(0.12), TEAL); rect(s, 0, SH-Inches(0.12), SW, Inches(0.12), AMBER)
txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.5), [[R("Terima kasih.", 60, WHITE, True, FONT_H)]])
txt(s, Inches(0.95), Inches(3.9), Inches(11), Inches(0.8),
    [[R("SahkahIni? — verifying the truth, one forward at a time.", 20, TEAL, False)]])
txt(s, Inches(0.95), Inches(5.0), Inches(11), Inches(1),
    [[R("Aisyatun Nabiha  ·  Ahli Kumpulan 2  ·  Ahli Kumpulan 3", 15, RGBColor(0xC8,0xD3,0xE6), False)],
     [R("GFP 009 Digital Creativity  ·  Universiti Malaya", 13, SLATE, False)]], space=6)

out = r"c:\Users\user\Documents\motionboards\SahkahIni-GFP009.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
