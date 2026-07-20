#!/usr/bin/env python3
"""Build the MicroMentor GFP 0009 deck (student version).

Voice: Aisyatun's — short confident sentences, benefit-led, British/Malaysian
spelling (utilise, personalise, minimise), a small tagline under the wordmark.
Design: Apple product-keynote feel. Helvetica Neue, generous air, one accent.
No long dashes anywhere. Hero imagery rendered with Nano Banana 2 (see
scripts/gen-micromentor.mjs) and full-bleed so it blends with the slide ground.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
PAPER   = RGBColor(0xEC, 0xEE, 0xF7)
PAPER2  = RGBColor(0xE3, 0xE6, 0xF3)
INK     = RGBColor(0x15, 0x18, 0x2B)
INKSOFT = RGBColor(0x4A, 0x4F, 0x6B)
CORAL   = RGBColor(0xFF, 0x5A, 0x47)
CORALDP = RGBColor(0xE8, 0x45, 0x2F)
HI      = RGBColor(0xFF, 0xD2, 0x3F)
INDIGO  = RGBColor(0x4B, 0x4F, 0xA6)
INDIGOS = RGBColor(0xE9, 0xEA, 0xFB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xC7, 0xCA, 0xDD)

FONT_H = "Helvetica Neue"
FONT_B = "Helvetica Neue"

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def img(s, name, l, t, w, h, back=False):
    pic = s.shapes.add_picture(os.path.join(ASSETS, name), l, t, width=w, height=h)
    if back:
        sp = pic._element; sp.getparent().remove(sp); s.shapes._spTree.insert(3, sp)
    return pic

def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf

def para(tf, text, size, color=INK, bold=False, first=False, space_before=0,
         space_after=6, align=PP_ALIGN.LEFT, font=FONT_B, spacing=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    if spacing: p.line_spacing = spacing
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    return p, r

def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def eyebrow(s, label):
    rect(s, Inches(0.6), Inches(0.55), Inches(0.14), Inches(0.42), fill=CORAL)
    tb, tf = box(s, Inches(0.86), Inches(0.5), Inches(8), Inches(0.55), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "GFP 0009 · DIGITAL CREATIVITY"
    r.font.size = Pt(12); r.font.bold = True; r.font.name = FONT_H
    r.font.color.rgb = INK
    _letterspace(r, 220)
    tb2, tf2 = box(s, Inches(0.86), Inches(0.9), Inches(9), Inches(0.4))
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(11); r2.font.bold = True; r2.font.name = FONT_H
    r2.font.color.rgb = INDIGO; _letterspace(r2, 160)

def footer(s, n, dark=False):
    c1 = RGBColor(0x9A,0x9F,0xBE) if dark else INKSOFT
    tb, tf = box(s, Inches(0.6), Inches(7.02), Inches(6), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "MicroMentor  ·  Beat the forgetting curve"
    r.font.size = Pt(9); r.font.name = FONT_B; r.font.color.rgb = c1
    tb2, tf2 = box(s, Inches(11.9), Inches(7.02), Inches(0.9), Inches(0.35))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"{n:02d} / 12"
    r2.font.size = Pt(9); r2.font.bold = True; r2.font.name = FONT_B
    r2.font.color.rgb = (HI if dark else INDIGO)

def _letterspace(run, val):
    rPr = run._r.get_or_add_rPr(); rPr.set('spc', str(val))

def _highlight(run, color=HI):
    rPr = run._r.get_or_add_rPr()
    h = rPr.makeelement(qn('a:highlight'), {})
    clr = h.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0], color[1], color[2])})
    h.append(clr); rPr.append(h)

def title_of(s, text, y=1.55, size=40, w=12.1, color=INK):
    tb, tf = box(s, Inches(0.6), Inches(y), Inches(w), Inches(1.3))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.name = FONT_H; r.font.color.rgb = color
    return tb

def bullet(tf, text, size=16, color=INK, bold=False, first=False, sub=False, space_after=10):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.line_spacing = 1.12
    p.level = 1 if sub else 0
    dot = p.add_run(); dot.text = ("·  " if sub else "●  ")
    dot.font.size = Pt(size if not sub else size-1)
    dot.font.color.rgb = (INDIGO if sub else CORAL); dot.font.name = FONT_B; dot.font.bold = True
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT_B; r.font.color.rgb = color
    return p, r

# =========================================================================
# SLIDE 1 — TITLE  (hero-light phone render on the right)
# =========================================================================
s = slide(PAPER)
img(s, "hero-light.png", Inches(8.05), Inches(-0.35), Inches(6.0), Inches(8.03), back=True)

rect(s, Inches(0.6), Inches(0.6), Inches(0.16), Inches(0.5), fill=CORAL)
tb, tf = box(s, Inches(0.9), Inches(0.55), Inches(7), Inches(0.6), MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "GFP 0009 · DIGITAL CREATIVITY"
r.font.size=Pt(13); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=INK; _letterspace(r,240)

# wordmark
tb, tf = box(s, Inches(0.6), Inches(1.75), Inches(7.6), Inches(1.4))
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text="Micro"; r1.font.size=Pt(70); r1.font.bold=True; r1.font.name=FONT_H; r1.font.color.rgb=INK
r2 = p.add_run(); r2.text="Mentor"; r2.font.size=Pt(70); r2.font.bold=True; r2.font.name=FONT_H; r2.font.color.rgb=CORALDP
_highlight(r2, HI)
r3 = p.add_run(); r3.text="."; r3.font.size=Pt(70); r3.font.bold=True; r3.font.name=FONT_H; r3.font.color.rgb=CORAL

# Aisyatun-style tagline under the wordmark
tb, tf = box(s, Inches(0.62), Inches(3.05), Inches(7), Inches(0.5))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "A LEARNING APP FOR STUDENTS"
r.font.size=Pt(13); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=INDIGO; _letterspace(r,200)

tb, tf = box(s, Inches(0.6), Inches(3.7), Inches(7.2), Inches(1.5))
p,_ = para(tf, "Turn a whole lecture into 3 minute lessons you actually remember.",
     25, INK, bold=True, first=True, spacing=1.12, space_after=8)
para(tf, "A student micro-learning app that beats cramming and the forgetting curve. Built on the AMPLIFY model.",
     14, INKSOFT, spacing=1.2)

# presenters
tb, tf = box(s, Inches(0.6), Inches(5.5), Inches(7.2), Inches(1.4))
para(tf, "PRESENTED BY", 11, INDIGO, bold=True, first=True, space_after=8)
for nm, mat, ghost in [("Aisyatun Nabiha binti Mohamed Husain","25004128",False),
                       ("[ Group member 2 ]","[ matric ]",True),
                       ("[ Group member 3 ]","[ matric ]",True)]:
    p = tf.add_paragraph(); p.space_after=Pt(4)
    d=p.add_run(); d.text="●  "; d.font.size=Pt(13); d.font.color.rgb=(GREY if ghost else CORAL); d.font.name=FONT_B
    rn=p.add_run(); rn.text=nm+"  "; rn.font.size=Pt(15); rn.font.bold=True; rn.font.name=FONT_B
    rn.font.color.rgb=(GREY if ghost else INK)
    rm=p.add_run(); rm.text=mat; rm.font.size=Pt(12); rm.font.name=FONT_B; rm.font.color.rgb=(GREY if ghost else INKSOFT)

# course meta, single line bottom-left
tb, tf = box(s, Inches(0.6), Inches(6.95), Inches(9), Inches(0.4))
p = tf.paragraphs[0]
r = p.add_run(); r.text="GFP 0009 Digital Creativity   |   Dr. Rafiza binti Abdul Razak   |   Universiti Malaya"
r.font.size=Pt(11); r.font.name=FONT_B; r.font.color.rgb=INKSOFT

# =========================================================================
# SLIDE 2 — THE PROBLEM
# =========================================================================
s = slide(PAPER); eyebrow(s, "THE PROBLEM"); footer(s,2)
title_of(s, "We cram. Then we forget.", size=42)
tb, tf = box(s, Inches(0.6), Inches(2.75), Inches(7.5), Inches(3))
bullet(tf, "Lectures dump everything at once. Working memory overloads.", 17, first=True, space_after=15)
bullet(tf, "Without review, we forget most of it within days (Ebbinghaus, 1885).", 17, space_after=15)
bullet(tf, "So we cram. High stress, and it fades right after.", 17, space_after=15)
rect(s, Inches(0.6), Inches(5.55), Inches(4.7), Inches(1.0), fill=WHITE, rounded=True, radius=0.14)
tb, tf = box(s, Inches(0.95), Inches(5.55), Inches(4.2), Inches(1.0), MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]
rb=p.add_run(); rb.text="50-80%  "; rb.font.size=Pt(32); rb.font.bold=True; rb.font.name=FONT_H; rb.font.color.rgb=CORALDP
rl=p.add_run(); rl.text="forgotten within days"; rl.font.size=Pt(13); rl.font.name=FONT_B; rl.font.color.rgb=INKSOFT
img(s, "spot-problem.png", Inches(8.55), Inches(2.15), Inches(3.7), Inches(4.63))

# =========================================================================
# SLIDE 3 — WHY IT MATTERS
# =========================================================================
s = slide(PAPER); eyebrow(s, "SIGNIFICANCE & JUSTIFICATION"); footer(s,3)
title_of(s, "Why this matters to us", size=42)
cols = [("Personal","Every one of us has pulled an all nighter, then forgotten it a week later. This is our own problem. Real, not hypothetical."),
        ("Scale","Thousands of UM students juggle heavy content across many courses. A tool that fits their phone habits helps at scale."),
        ("Backed by research","Micro-learning is not a guess. It is supported by cognitive science and by AMPLIFY, a framework validated in Malaysian higher education."),
        ("On theme","GFP 0009 asks us to use creativity to design appropriate technology for a real problem. Learning is exactly that.")]
for i,(h,b) in enumerate(cols):
    cx = Inches(0.6 + (i%2)*6.15); cy = Inches(2.5 + (i//2)*2.05)
    rect(s, cx, cy, Inches(5.9), Inches(1.8), fill=WHITE, rounded=True, radius=0.08)
    rect(s, cx, cy, Inches(0.14), Inches(1.8), fill=CORAL)
    tb, tf = box(s, cx+Inches(0.4), cy+Inches(0.25), Inches(5.3), Inches(1.4))
    para(tf, h, 18, INK, bold=True, first=True, space_after=6)
    para(tf, b, 13.5, INKSOFT, spacing=1.18)

# =========================================================================
# SLIDE 4 — TECHNOLOGY REVIEW
# =========================================================================
s = slide(PAPER); eyebrow(s, "USE OF TECHNOLOGY IN THIS AREA"); footer(s,4)
title_of(s, "How students learn, then and now", size=34)
rows = [("Then","Textbooks, printed notes, classroom only","Not portable, passive, forgotten fast"),
        ("E-learning","LMS like Spectrum and Moodle, PDF slides, recorded lectures","Long courses on a screen, same overload"),
        ("Video","YouTube, recorded webinars","Passive, no feedback, easy to zone out"),
        ("Now","Micro-learning apps like Duolingo, Anki, Quizlet","Short, spaced, gamified, but generic and not tied to your course")]
tw=Inches(12.1); tx=Inches(0.6); ty=Inches(2.35); rh=Inches(0.95)
hdr = ["ERA","TECHNOLOGY USED","STRENGTH AND LIMITATION"]
cw=[Inches(2.0),Inches(5.0),Inches(5.1)]
hx=tx
rect(s, tx, ty, tw, Inches(0.55), fill=INK)
for j,ht in enumerate(hdr):
    tb, tf = box(s, hx+Inches(0.2), ty, cw[j]-Inches(0.3), Inches(0.55), MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=ht; r.font.size=Pt(11); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=WHITE; _letterspace(r,120)
    hx+=cw[j]
for i,row in enumerate(rows):
    ry=ty+Inches(0.55)+ i*Emu(int(rh))
    is_now = row[0]=="Now"
    rect(s, tx, ry, tw, rh, fill=(INDIGOS if is_now else WHITE), line=PAPER2, line_w=1)
    hx=tx
    for j,cell in enumerate(row):
        tb, tf = box(s, hx+Inches(0.2), ry, cw[j]-Inches(0.3), rh, MSO_ANCHOR.MIDDLE)
        p=tf.paragraphs[0]; r=p.add_run(); r.text=cell
        r.font.size=Pt(12.5 if j else 13); r.font.name=FONT_B
        r.font.bold = (j==0) or is_now
        r.font.color.rgb = (CORALDP if (j==0 and is_now) else INK if j==0 else INKSOFT)
        hx+=cw[j]
tb, tf = box(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5))
para(tf, "The gap. Existing apps are not tied to your lectures. MicroMentor turns your own course material into micro-lessons.",
     13, INDIGO, bold=True, first=True)

# =========================================================================
# SLIDE 5 — OUR SOLUTION  (hero-dark full-bleed, text on the left)
# =========================================================================
s = slide(INK); footer(s,5, dark=True)
img(s, "hero-dark.png", 0, 0, SW, SH, back=True)
rect(s, Inches(0.6), Inches(0.6), Inches(0.16), Inches(0.5), fill=CORAL)
tb, tf = box(s, Inches(0.9), Inches(0.55), Inches(7), Inches(0.6), MSO_ANCHOR.MIDDLE)
r=tf.paragraphs[0].add_run(); r.text="OUR SOLUTION"; r.font.size=Pt(12); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=HI; _letterspace(r,220)
tb, tf = box(s, Inches(0.6), Inches(1.95), Inches(7.2), Inches(2.6))
p=tf.paragraphs[0]; p.line_spacing=1.1
r=p.add_run(); r.text="MicroMentor turns your "; r.font.size=Pt(34); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=WHITE
r=p.add_run(); r.text="lecture notes"; r.font.size=Pt(34); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=HI
r=p.add_run(); r.text=" into short, personalised lessons, with quizzes and spaced reminders so it "; r.font.size=Pt(34); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=WHITE
r=p.add_run(); r.text="sticks."; r.font.size=Pt(34); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=CORAL
feats=[("3-5 min lessons","One idea at a time, on your phone."),
       ("Smart review","Reminds you right before you would forget."),
       ("Learn together","A class leaderboard and a Q&A wall.")]
for i,(h,b) in enumerate(feats):
    yy=Inches(5.05+i*0.62)
    rect(s, Inches(0.6), yy, Inches(0.12), Inches(0.5), fill=CORAL)
    tb, tf = box(s, Inches(0.9), yy, Inches(6.6), Inches(0.55), MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]
    rh_=p.add_run(); rh_.text=h+"   "; rh_.font.size=Pt(15); rh_.font.bold=True; rh_.font.name=FONT_B; rh_.font.color.rgb=HI
    rb_=p.add_run(); rb_.text=b; rb_.font.size=Pt(13.5); rb_.font.name=FONT_B; rb_.font.color.rgb=RGBColor(0xC9,0xCC,0xE0)

# =========================================================================
# SLIDE 6 — AMPLIFY APPLIED (key slide)
# AMPLIFY is the model's NAME (Advancing Microlearning Practices for Enhanced
# Learning Experiences). Its seven components are a plain list, NOT letter
# initials, so we number them 1-7 rather than mapping them to A/M/P/L/I/F/Y.
# =========================================================================
s = slide(PAPER); eyebrow(s, "APPLYING THE AMPLIFY MODEL"); footer(s,6)
title_of(s, "AMPLIFY, applied to MicroMentor", size=34, y=1.5)
tb, tf = box(s, Inches(0.6), Inches(2.12), Inches(12.1), Inches(0.4))
para(tf,"AMPLIFY is a microlearning model by Nurul Fitriah Alias (UM, 2024). Its seven components become our seven design decisions.",
     12.5, INKSOFT, italic=True, first=True, space_after=0)
amp=[("Clear learning objectives","One goal per lesson, like explain the light reaction."),
     ("Personalised content","Adapts to your course and your weak spots."),
     ("Continuous feedback","Instant answers plus spaced review reminders."),
     ("Efficient content delivery","3 to 5 minutes, mobile first, works offline."),
     ("Interactive experiences","Tap, drag and scenario questions. Never passive."),
     ("Social interaction","A class leaderboard and a peer Q&A wall."),
     ("Seamless curriculum integration","Lessons map to your real syllabus.")]
ty=Inches(2.62); rh=Inches(0.585)
for i,(h,b) in enumerate(amp):
    ry=ty+ i*Emu(int(rh))
    rect(s, Inches(0.6), ry+Inches(0.02), Inches(0.5), Inches(0.5), fill=CORAL, rounded=True, radius=0.25)
    tb, tf = box(s, Inches(0.6), ry+Inches(0.02), Inches(0.5), Inches(0.5), MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=str(i+1)
    r.font.size=Pt(20); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=WHITE
    tb, tf = box(s, Inches(1.3), ry, Inches(4.7), Inches(0.55), MSO_ANCHOR.MIDDLE)
    para(tf, h, 14.5, INK, bold=True, first=True, space_after=0)
    tb, tf = box(s, Inches(6.1), ry, Inches(6.6), Inches(0.55), MSO_ANCHOR.MIDDLE)
    para(tf, b, 13.5, INKSOFT, first=True, space_after=0)

# =========================================================================
# SLIDE 7 — STORYBOARD
# =========================================================================
s = slide(PAPER); eyebrow(s, "DESIGNING THE SOLUTION · STORYBOARD"); footer(s,7)
title_of(s, "Aina's story in six frames", size=34)
frames=[("1","The problem","Aina faces 120 slides the night before her exam. Panic."),
        ("2","Discovery","She opens MicroMentor. It says do three short lessons today."),
        ("3","Learning","A three minute lesson, one clear goal, tap to answer."),
        ("4","Feedback","Instant tick and a reason. A reminder brings it back tomorrow."),
        ("5","Together","She checks the class leaderboard and asks a question on the wall."),
        ("6","Outcome","Exam week arrives. She feels calm, confident, and actually remembers.")]
for i,(n,h,b) in enumerate(frames):
    cx=Inches(0.6+(i%3)*4.05); cy=Inches(2.35+(i//3)*2.15)
    rect(s, cx, cy, Inches(3.8), Inches(1.95), fill=WHITE, rounded=True, radius=0.06)
    rect(s, cx+Inches(0.3), cy+Inches(0.3), Inches(0.55), Inches(0.55), fill=INDIGOS, rounded=True, radius=0.25)
    tb, tf = box(s, cx+Inches(0.3), cy+Inches(0.3), Inches(0.55), Inches(0.55), MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=n
    r.font.size=Pt(20); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=INDIGO
    tb, tf = box(s, cx+Inches(1.05), cy+Inches(0.32), Inches(2.5), Inches(0.55), MSO_ANCHOR.MIDDLE)
    para(tf, h, 16, INK, bold=True, first=True, space_after=0)
    tb, tf = box(s, cx+Inches(0.32), cy+Inches(1.0), Inches(3.2), Inches(0.9))
    para(tf, b, 12.5, INKSOFT, first=True, spacing=1.15)

# =========================================================================
# SLIDE 8 — PROTOTYPE  (phones-trio render as the hero)
# =========================================================================
s = slide(PAPER); eyebrow(s, "PROTOTYPE"); footer(s,8)
title_of(s, "The clickable prototype", size=36, y=1.5)
tb, tf = box(s, Inches(0.6), Inches(2.45), Inches(5.5), Inches(4.2))
bullet(tf,"Built in Figma. Seven core screens, fully clickable.",15.5,first=True)
bullet(tf,"Onboarding. Pick your course and level, which feeds M.",14.5,sub=True)
bullet(tf,"Today's lessons home, two or three bite sized cards (A, L).",14.5,sub=True)
bullet(tf,"Lesson screen, a short visual plus a tap question (I).",14.5,sub=True)
bullet(tf,"Instant feedback, the reason, and your streak (P).",14.5,sub=True)
bullet(tf,"Class wall and leaderboard (F).",14.5,sub=True)
bullet(tf,"Syllabus playlists mapped to the course (Y).",14.5,sub=True)
# trio render on the right (16:9), placed to bleed the right edge
img(s, "phones-trio.png", Inches(6.05), Inches(2.55), Inches(7.6), Inches(4.24))
tb, tf = box(s, Inches(6.4), Inches(6.5), Inches(6.5), Inches(0.4))
para(tf,"Home · progress · class wall",11,INKSOFT,italic=True,first=True,space_after=0)

# =========================================================================
# SLIDE 9 — TESTING & ITERATION
# =========================================================================
s = slide(PAPER); eyebrow(s, "TESTING & EVALUATION"); footer(s,9)
title_of(s, "We tested with peers, then changed four things", size=32)
pairs=[("Feedback did not say why an answer was wrong.","We added a one line explanation to every result."),
       ("Making a lesson took too many steps.","We cut the create lesson flow from six steps to three."),
       ("The leaderboard felt stressful for some.","We made it opt in and class only."),
       ("People wanted a language choice.","We added a Bahasa Melayu and English toggle.")]
tb, tf = box(s, Inches(0.6), Inches(2.3), Inches(5.9), Inches(0.5))
para(tf,"WHAT USERS SAID",12,CORALDP,bold=True,first=True,space_after=0)
tb, tf = box(s, Inches(6.9), Inches(2.3), Inches(5.9), Inches(0.5))
para(tf,"WHAT WE CHANGED",12,INDIGO,bold=True,first=True,space_after=0)
for i,(a,b) in enumerate(pairs):
    yy=Inches(2.85+i*1.0)
    rect(s, Inches(0.6), yy, Inches(5.5), Inches(0.85), fill=WHITE, rounded=True, radius=0.1)
    tb, tf = box(s, Inches(0.85), yy, Inches(5.0), Inches(0.85), MSO_ANCHOR.MIDDLE)
    para(tf,a,13.5,INK,first=True,space_after=0,spacing=1.1)
    ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.25), yy+Inches(0.27), Inches(0.5), Inches(0.3))
    ar.fill.solid(); ar.fill.fore_color.rgb=CORAL; ar.line.fill.background(); ar.shadow.inherit=False
    rect(s, Inches(6.9), yy, Inches(5.5), Inches(0.85), fill=INDIGOS, rounded=True, radius=0.1)
    tb, tf = box(s, Inches(7.15), yy, Inches(5.0), Inches(0.85), MSO_ANCHOR.MIDDLE)
    para(tf,b,13.5,INK,bold=True,first=True,space_after=0,spacing=1.1)

# =========================================================================
# SLIDE 10 — PROJECT MANAGEMENT
# =========================================================================
s = slide(PAPER); eyebrow(s, "PROJECT MANAGEMENT"); footer(s,10)
title_of(s, "Timeline, sustainability and commercialisation", size=30)
tb, tf = box(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.4))
para(tf,"14 WEEK TIMELINE, ALIGNED TO THE COURSE SCHEDULE",12,CORALDP,bold=True,first=True,space_after=0)
phases=[("Wk 1-4","Research\nthe problem"),("Wk 5","Map\nAMPLIFY"),("Wk 6-7","Ideate,\nAssign 1"),
        ("Wk 8-9","Prototype,\nstoryboard"),("Wk 10-11","Usability\ntest"),("Wk 12-14","Iterate,\nAssign 2")]
rect(s, Inches(0.7), Inches(2.95), Inches(11.9), Inches(0.05), fill=GREY)
for i,(w,lab) in enumerate(phases):
    cx=Inches(0.6+i*2.03)
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL, cx+Inches(0.85), Inches(2.85), Inches(0.24), Inches(0.24))
    dot.fill.solid(); dot.fill.fore_color.rgb=CORAL; dot.line.fill.background(); dot.shadow.inherit=False
    tb, tf = box(s, cx, Inches(3.2), Inches(1.9), Inches(0.35))
    para(tf,w,12,INK,bold=True,first=True,align=PP_ALIGN.CENTER,space_after=0)
    tb, tf = box(s, cx, Inches(3.55), Inches(1.9), Inches(0.7))
    for j,ln in enumerate(lab.split("\n")):
        para(tf,ln,11,INKSOFT,first=(j==0),align=PP_ALIGN.CENTER,space_after=0,spacing=1.0)
cards=[("Sustainability","Free tools only, like Figma, Canva and the UM e-portfolio. No cost to the course. Going digital means no printed notes."),
       ("Commercialisation","Freemium. Free for the basics, a low fee for premium. Sell ready made lesson packs by course. Bilingual in BM and English."),
       ("Creativity extras","Bilingual with audio, an offline mode, and gentle analytics that coach you rather than rank you.")]
for i,(h,b) in enumerate(cards):
    cx=Inches(0.6+i*4.05)
    rect(s, cx, Inches(4.7), Inches(3.75), Inches(2.0), fill=WHITE, rounded=True, radius=0.07)
    rect(s, cx, Inches(4.7), Inches(3.75), Inches(0.12), fill=CORAL)
    tb, tf = box(s, cx+Inches(0.35), Inches(5.0), Inches(3.1), Inches(1.6))
    para(tf,h,16,INK,bold=True,first=True,space_after=6)
    para(tf,b,12.5,INKSOFT,spacing=1.18)

# =========================================================================
# SLIDE 11 — FUTURE & CONCLUSION
# =========================================================================
s = slide(PAPER); eyebrow(s, "FUTURE USE & CONCLUSION"); footer(s,11)
title_of(s, "Where it is heading", size=40)
tb, tf = box(s, Inches(0.6), Inches(2.4), Inches(6.2), Inches(4))
bullet(tf,"AI turns any PDF of notes into lessons and quizzes in minutes.",16,first=True)
bullet(tf,"It lives inside WhatsApp and Teams, so you learn where you already are.",16)
bullet(tf,"Micro-credentials and badges that follow you between courses.",16)
bullet(tf,"We design it responsibly. No notification fatigue, private data, always bilingual.",16)
rect(s, Inches(7.1), Inches(2.4), Inches(5.6), Inches(3.9), fill=INK, rounded=True, radius=0.05)
tb, tf = box(s, Inches(7.5), Inches(2.8), Inches(4.9), Inches(3.2))
para(tf,"THE TAKEAWAY",11,HI,bold=True,first=True,space_after=12)
para(tf,"Creativity is most powerful when it is aimed at a real problem and structured by a sound model.",19,WHITE,bold=True,spacing=1.2,space_after=12)
para(tf,"MicroMentor applies AMPLIFY to make learning short, personal and sticky. Appropriate technology for students.",14,RGBColor(0xC9,0xCC,0xE0),spacing=1.25)

# =========================================================================
# SLIDE 12 — THANK YOU  (hero-dark bookend)
# =========================================================================
s = slide(INK); footer(s,12, dark=True)
img(s, "hero-dark.png", 0, 0, SW, SH, back=True)
tb, tf = box(s, Inches(0.9), Inches(2.4), Inches(7), Inches(2))
p=tf.paragraphs[0]
r=p.add_run(); r.text="Terima kasih"; r.font.size=Pt(60); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=WHITE
r=p.add_run(); r.text="."; r.font.size=Pt(60); r.font.bold=True; r.font.name=FONT_H; r.font.color.rgb=CORAL
tb, tf = box(s, Inches(0.9), Inches(3.8), Inches(7), Inches(0.9))
p=tf.paragraphs[0]
r=p.add_run(); r.text="Questions? Let us "; r.font.size=Pt(23); r.font.name=FONT_B; r.font.color.rgb=RGBColor(0xC9,0xCC,0xE0)
r=p.add_run(); r.text="MicroMentor"; r.font.size=Pt(23); r.font.bold=True; r.font.name=FONT_B; r.font.color.rgb=HI
r=p.add_run(); r.text=" your notes."; r.font.size=Pt(23); r.font.name=FONT_B; r.font.color.rgb=RGBColor(0xC9,0xCC,0xE0)
tb, tf = box(s, Inches(0.9), Inches(5.5), Inches(8), Inches(0.6))
para(tf,"Aisyatun Nabiha binti Mohamed Husain · 25004128    |    GFP 0009 · Dr. Rafiza binti Abdul Razak",12.5,RGBColor(0x9A,0x9F,0xBE),first=True)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "MicroMentor-GFP0009.pptx")
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
