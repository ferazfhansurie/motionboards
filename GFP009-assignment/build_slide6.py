#!/usr/bin/env python3
"""Standalone corrected SLIDE 6 only, so it can be dropped into an already
adjusted deck without regenerating (and clobbering) the other slides.
AMPLIFY is the model's NAME, not the 7 components -> number them 1-7."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PAPER=RGBColor(0xEC,0xEE,0xF7); INK=RGBColor(0x15,0x18,0x2B); INKSOFT=RGBColor(0x4A,0x4F,0x6B)
CORAL=RGBColor(0xFF,0x5A,0x47); INDIGO=RGBColor(0x4B,0x4F,0xA6); WHITE=RGBColor(0xFF,0xFF,0xFF)
FONT="Helvetica Neue"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
s=prs.slides.add_slide(prs.slide_layouts[6])
def bg():
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=PAPER
    r.line.fill.background(); r.shadow.inherit=False
    sp=r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
bg()
def box(l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0; return tb,tf
def para(tf,text,size,color=INK,bold=False,first=False,space_after=6,italic=False,align=PP_ALIGN.LEFT):
    p=tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(space_after)
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.name=FONT; r.font.color.rgb=color; return p,r
def ls(run,val): run._r.get_or_add_rPr().set('spc',str(val))
def rrect(l,t,w,h,fill,rounded=True,radius=0.25):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,l,t,w,h)
    if rounded:
        try: shp.adjustments[0]=radius
        except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb=fill; shp.line.fill.background(); shp.shadow.inherit=False; return shp

# eyebrow
rrect(Inches(0.6),Inches(0.55),Inches(0.14),Inches(0.42),CORAL,rounded=False)
tb,tf=box(Inches(0.86),Inches(0.5),Inches(8),Inches(0.55),MSO_ANCHOR.MIDDLE)
_,r=para(tf,"GFP 0009 · DIGITAL CREATIVITY",12,INK,bold=True,first=True,space_after=0); ls(r,220)
tb,tf=box(Inches(0.86),Inches(0.9),Inches(9),Inches(0.4))
_,r=para(tf,"APPLYING THE AMPLIFY MODEL",11,INDIGO,bold=True,first=True,space_after=0); ls(r,160)
# title
tb,tf=box(Inches(0.6),Inches(1.5),Inches(12.1),Inches(0.9))
para(tf,"AMPLIFY, applied to MicroMentor",34,INK,bold=True,first=True)
# caption
tb,tf=box(Inches(0.6),Inches(2.12),Inches(12.1),Inches(0.4))
para(tf,"AMPLIFY is a microlearning model by Nurul Fitriah Alias (UM, 2024). Its seven components become our seven design decisions.",12.5,INKSOFT,italic=True,first=True,space_after=0)
# 7 numbered components
amp=[("Clear learning objectives","One goal per lesson, like explain the light reaction."),
     ("Personalised content","Adapts to your course and your weak spots."),
     ("Continuous feedback","Instant answers plus spaced review reminders."),
     ("Efficient content delivery","3 to 5 minutes, mobile first, works offline."),
     ("Interactive experiences","Tap, drag and scenario questions. Never passive."),
     ("Social interaction","A class leaderboard and a peer Q&A wall."),
     ("Seamless curriculum integration","Lessons map to your real syllabus.")]
ty=Inches(2.62); rh=Inches(0.585)
for i,(h,b) in enumerate(amp):
    ry=ty+i*Emu(int(rh))
    rrect(Inches(0.6),ry+Inches(0.02),Inches(0.5),Inches(0.5),CORAL,radius=0.25)
    tb,tf=box(Inches(0.6),ry+Inches(0.02),Inches(0.5),Inches(0.5),MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; rn=p.add_run(); rn.text=str(i+1)
    rn.font.size=Pt(20); rn.font.bold=True; rn.font.name=FONT; rn.font.color.rgb=WHITE
    tb,tf=box(Inches(1.3),ry,Inches(4.7),Inches(0.55),MSO_ANCHOR.MIDDLE)
    para(tf,h,14.5,INK,bold=True,first=True,space_after=0)
    tb,tf=box(Inches(6.1),ry,Inches(6.6),Inches(0.55),MSO_ANCHOR.MIDDLE)
    para(tf,b,13.5,INKSOFT,first=True,space_after=0)

out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"MicroMentor-Slide6-fixed.pptx")
prs.save(out); print("saved",out)
